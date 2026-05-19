"""
Gerador de apresentação PPT profissional — LLE Índices.

Slides ricos com gráficos reais (matplotlib fallback sem kaleido),
cards de KPI detalhados e análise profunda para apresentação à diretoria.
"""
from __future__ import annotations
import io
from typing import Optional

# Cores LLE
AZUL     = "041747"
AMARELO  = "FAC318"
VERDE    = "0F8C3B"
VERMELHO = "DC3545"
BRANCO   = "FFFFFF"
CINZA_F  = "F2F4F8"
CINZA_M  = "6C757D"
AZUL_V   = "0071FE"
CINZA_L  = "E8EDF5"


# ============================================================
# ENTRYPOINTS PÚBLICOS
# ============================================================

def gerar_ppt_cobranca(dados: dict, mes_label: str) -> bytes:
    prs = _nova_apresentacao()
    _slide_capa(prs, "Relatório de Cobrança", mes_label, "Gestão de Cobrança · Grupo LLE")
    _slide_kpis_cobranca(prs, dados, mes_label)
    _slide_acordos_detalhado(prs, dados, mes_label)
    _slide_baixas_detalhado(prs, dados, mes_label)
    _slide_performance_detalhado(prs, dados, mes_label)
    _slide_destaques_texto(prs, _analise_cobranca(dados), "ANÁLISE DO PERÍODO — COBRANÇA", mes_label)
    _slide_encerramento(prs, "COBRANÇA")
    return _salvar(prs)


def gerar_ppt_credito(dados: dict, mes_label: str) -> bytes:
    prs = _nova_apresentacao()
    _slide_capa(prs, "Relatório de Crédito", mes_label, "Análise de Crédito · Grupo LLE")
    _slide_kpis_credito(prs, dados, mes_label)
    _slide_liberacoes_detalhado(prs, dados, mes_label)
    _slide_analistas_detalhado(prs, dados, mes_label)
    _slide_destaques_texto(prs, _analise_credito(dados), "ANÁLISE DO PERÍODO — CRÉDITO", mes_label)
    _slide_encerramento(prs, "CRÉDITO")
    return _salvar(prs)


def gerar_ppt_reanalises(dados: dict, mes_label: str) -> bytes:
    prs = _nova_apresentacao()
    _slide_capa(prs, "Reanálises de Limite de Crédito", mes_label, "Análise de Crédito · Grupo LLE")
    _slide_kpis_reanalises(prs, dados, mes_label)
    _slide_reanalises_graficos(prs, dados, mes_label)
    _slide_top_reanalises(prs, dados, mes_label)
    _slide_destaques_texto(prs, _analise_reanalises(dados), "ANÁLISE DE REANÁLISES", mes_label)
    _slide_encerramento(prs, "REANÁLISES DE LIMITE")
    return _salvar(prs)


def gerar_ppt_geral(dados_cobranca: list[dict], dados_credito: list[dict]) -> bytes:
    """PPT Geral — análise profunda de todos os índices de Crédito e Cobrança."""
    prs = _nova_apresentacao()
    _slide_capa(prs, "Relatório Geral — Crédito & Cobrança",
                "Análise Consolidada", "Gestão Financeira · Grupo LLE")
    _slide_sumario_executivo(prs, dados_cobranca, dados_credito)
    _slide_evolucao_cobranca_geral(prs, dados_cobranca)
    _slide_analise_cobradores_geral(prs, dados_cobranca)
    _slide_evolucao_credito_geral(prs, dados_credito)
    _slide_analise_analistas_geral(prs, dados_credito)
    _slide_tendencias_geral(prs, dados_cobranca, dados_credito)
    _slide_analise_profunda_texto(prs, dados_cobranca, dados_credito)
    _slide_alertas_recomendacoes(prs, dados_cobranca, dados_credito)
    _slide_encerramento(prs, "GESTÃO FINANCEIRA")
    return _salvar(prs)


# ============================================================
# SLIDES — COBRANÇA INDIVIDUAL
# ============================================================

def _slide_kpis_cobranca(prs, dados: dict, periodo: str):
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "RESUMO EXECUTIVO — COBRANÇA", periodo)

    from pptx.util import Inches
    ac = dados.get("acordos", {})
    bx = dados.get("baixas", {})
    pf = dados.get("performance", {})

    kpis = [
        ("Total de Acordos", str(ac.get("total_acordos", 0)), "realizados no período", AZUL_V),
        ("Valor Total Acordado", _brl(ac.get("valor_total", 0)), "soma dos acordos ativos", VERDE),
        ("Ticket Médio", _brl(ac.get("ticket_medio", 0)), "valor médio por acordo", AZUL_V),
        ("Total Recebido", _brl(bx.get("total_recebido", 0)), "cobrança baixada no período", VERDE),
        ("Taxa de Cancelamento", _pct(ac.get("pct_cancelamento", 0)), "acordos cancelados", VERMELHO),
        ("Aderência à Meta", _pct(pf.get("aderencia_media_pct", 0)), "meta de 2h/dia por cobrador", AZUL_V),
    ]
    _grid_kpis(slide, kpis, Inches(1.5))


def _slide_acordos_detalhado(prs, dados: dict, periodo: str):
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "ÍNDICES DE ACORDO — DETALHAMENTO", periodo)

    df = dados.get("df_acordos")
    if df is None or (hasattr(df, "empty") and df.empty):
        _texto_sem_dados(slide)
        return

    import pandas as pd
    df = pd.DataFrame(df) if isinstance(df, list) else df

    col_l, col_r = Inches(0.3), Inches(6.8)
    w, h = Inches(6.2), Inches(5.2)

    # Gráfico 1: Valor acordado por cobrador
    if "negociador" in df.columns and "valor_total" in df.columns:
        df_at = df[~df.get("cancelado", pd.Series([False]*len(df)))] if "cancelado" in df.columns else df
        por_neg = df_at.groupby("negociador")["valor_total"].sum().reset_index().sort_values("valor_total")
        fig1 = _fig_barras_h(
            por_neg["negociador"].tolist(),
            por_neg["valor_total"].tolist(),
            "Valor Acordado por Negociador (R$)", AZUL_V
        )
        _inserir_grafico(slide, fig1, col_l, Inches(1.5), w, h)

    # Gráfico 2: Forma de pagamento
    if "forma_pagto" in df.columns:
        df_at = df[~df["cancelado"]] if "cancelado" in df.columns else df
        por_forma = df_at["forma_pagto"].value_counts()
        fig2 = _fig_pizza(
            por_forma.index.tolist(),
            por_forma.values.tolist(),
            "Distribuição por Forma de Pagamento"
        )
        _inserir_grafico(slide, fig2, col_r, Inches(1.5), w, h)


def _slide_baixas_detalhado(prs, dados: dict, periodo: str):
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "ÍNDICES DE COBRANÇA — BAIXAS", periodo)

    df = dados.get("df_baixas")
    if df is None or (hasattr(df, "empty") and df.empty):
        _texto_sem_dados(slide)
        return

    import pandas as pd
    df = pd.DataFrame(df) if isinstance(df, list) else df

    col_l, col_r = Inches(0.3), Inches(6.8)
    w, h = Inches(6.2), Inches(5.2)

    # Gráfico 1: Recebido por cobrador
    if "cobrador" in df.columns and "vlr_liquido" in df.columns:
        por_cob = df.groupby("cobrador")["vlr_liquido"].sum().reset_index().sort_values("vlr_liquido")
        fig1 = _fig_barras_h(
            por_cob["cobrador"].tolist(),
            por_cob["vlr_liquido"].tolist(),
            "Total Recebido por Cobrador (R$)", VERDE
        )
        _inserir_grafico(slide, fig1, col_l, Inches(1.5), w, h)

    # Gráfico 2: Aging
    if "faixa_aging" in df.columns and "vlr_liquido" in df.columns:
        faixas = ["0-30", "31-60", "60-91", "91-180", "181-360", "+360"]
        cores_aging = [VERDE, "F5A623", "FF8C00", VERMELHO, "8B0000", "4A0000"]
        aging = df.groupby("faixa_aging")["vlr_liquido"].sum().reindex(faixas, fill_value=0)
        fig2 = _fig_barras_v_cores(
            faixas,
            aging.tolist(),
            "Aging — Valor por Faixa de Atraso (R$)",
            cores_aging
        )
        _inserir_grafico(slide, fig2, col_r, Inches(1.5), w, h)


def _slide_performance_detalhado(prs, dados: dict, periodo: str):
    from pptx.util import Inches, Pt
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "PERFORMANCE DOS COBRADORES", periodo)

    df = dados.get("df_performance")
    if df is None or (hasattr(df, "empty") and df.empty):
        _texto_sem_dados(slide)
        return

    import pandas as pd
    df = pd.DataFrame(df) if isinstance(df, list) else df
    df_ind = df[~df["cobrador"].str.upper().str.contains("EQUIPE|TOTAL", na=False)].copy() if "cobrador" in df.columns else df

    col_l, col_r = Inches(0.3), Inches(6.8)
    w, h = Inches(6.2), Inches(3.0)

    # Gráfico: TTO vs Meta
    if "cobrador" in df_ind.columns and "tempo_medio_diario_h" in df_ind.columns:
        cores = [VERDE if v >= 2.0 else VERMELHO for v in df_ind["tempo_medio_diario_h"]]
        fig1 = _fig_barras_v_cores_linha(
            df_ind["cobrador"].tolist(),
            df_ind["tempo_medio_diario_h"].tolist(),
            cores, 2.0, "Tempo Médio Diário vs Meta (2h)"
        )
        _inserir_grafico(slide, fig1, col_l, Inches(1.5), w, h)

    # Gráfico: Valor total acordos
    if "cobrador" in df_ind.columns and "total_acordos_valor" in df_ind.columns:
        df_sort = df_ind.sort_values("total_acordos_valor", ascending=True)
        fig2 = _fig_barras_h(
            df_sort["cobrador"].tolist(),
            df_sort["total_acordos_valor"].tolist(),
            "Valor Total de Acordos (R$)", AZUL_V
        )
        _inserir_grafico(slide, fig2, col_r, Inches(1.5), w, h)

    # Tabela de performance
    y = Inches(4.65)
    headers = ["Cobrador", "TTO Médio", "Aderência", "Ocorr./dia", "Acordos/dia", "Total R$"]
    larguras = [Inches(2.2), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(2.0)]
    _tabela_dados(slide, df_ind, headers,
                  ["cobrador", "tempo_medio_diario_h", "pct_aderencia",
                   "ocorrencias_media_dia", "acordos_media_dia", "total_acordos_valor"],
                  larguras, y, formatos={1: "h", 2: "%", 5: "R$"})


# ============================================================
# SLIDES — CRÉDITO INDIVIDUAL
# ============================================================

def _slide_kpis_credito(prs, dados: dict, periodo: str):
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "RESUMO EXECUTIVO — CRÉDITO", periodo)

    lib = dados.get("liberacoes", {})
    lim = dados.get("limites", {})

    kpis = [
        ("Total de Pedidos", f"{lib.get('total', 0):,}".replace(",", "."), "processados no período", AZUL_V),
        ("Passaram Direto", _pct(lib.get("pct_direto", 0)), f"{lib.get('qtd_direto', 0):,} pedidos automáticos".replace(",", "."), VERDE),
        ("Liberados (Analista)", _pct(lib.get("pct_liberados", 0)), f"{lib.get('qtd_liberados', 0):,} pedidos analisados".replace(",", "."), AZUL_V),
        ("Negados", _pct(lib.get("pct_negados", 0)), f"{lib.get('qtd_negados', 0):,} pedidos reprovados".replace(",", "."), VERMELHO),
        ("Valor Aprovado", _brl(lib.get("vlr_aprovado", 0)), "direto + analista", VERDE),
        ("Reanálises Limite", str(lim.get("total", 0)), f"{_brl(lim.get('total_variacao', 0))} concedido", AMARELO),
    ]
    _grid_kpis(slide, kpis, Inches(1.5))


def _slide_liberacoes_detalhado(prs, dados: dict, periodo: str):
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "FLUXO DE PEDIDOS DE CRÉDITO", periodo)

    df = dados.get("df_liberacoes")
    if df is None or (hasattr(df, "empty") and df.empty):
        _texto_sem_dados(slide)
        return

    import pandas as pd
    df = pd.DataFrame(df) if isinstance(df, list) else df

    col_l, col_r = Inches(0.3), Inches(6.8)
    w, h = Inches(6.2), Inches(5.2)

    total = len(df)
    qtd_d = len(df[df["tipo"] == "DIRETO"]) if "tipo" in df.columns else 0
    qtd_l = len(df[df["tipo"] == "LIBERADO"]) if "tipo" in df.columns else 0
    qtd_n = len(df[df["tipo"] == "NEGADO"]) if "tipo" in df.columns else 0

    fig1 = _fig_pizza(
        ["Passaram Direto", "Liberados (Analista)", "Negados"],
        [qtd_d, qtd_l, qtd_n],
        f"Distribuição de {total:,} Pedidos".replace(",", ".")
    )
    _inserir_grafico(slide, fig1, col_l, Inches(1.5), w, h)

    if "vlr_pedido" in df.columns and "tipo" in df.columns:
        vlr_d = float(df[df["tipo"] == "DIRETO"]["vlr_pedido"].sum())
        vlr_l = float(df[df["tipo"] == "LIBERADO"]["vlr_pedido"].sum())
        vlr_n = float(df[df["tipo"] == "NEGADO"]["vlr_pedido"].sum())
        fig2 = _fig_barras_v_cores(
            ["Direto", "Liberado (Analista)", "Negado"],
            [vlr_d, vlr_l, vlr_n],
            "Valor Total por Tipo de Pedido (R$)",
            [VERDE, AZUL_V, VERMELHO]
        )
        _inserir_grafico(slide, fig2, col_r, Inches(1.5), w, h)


def _slide_analistas_detalhado(prs, dados: dict, periodo: str):
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "RANKING DE ANALISTAS — LIBERAÇÕES", periodo)

    df = dados.get("df_liberacoes")
    if df is None or (hasattr(df, "empty") and df.empty):
        _texto_sem_dados(slide)
        return

    import pandas as pd
    df = pd.DataFrame(df) if isinstance(df, list) else df

    df_lib = df[df["tipo"] == "LIBERADO"].copy() if "tipo" in df.columns else df
    if "analista" not in df_lib.columns or df_lib["analista"].isna().all():
        _texto_sem_dados(slide)
        return

    ranking = (
        df_lib.dropna(subset=["analista"])
        .groupby("analista")
        .agg(qtd=("vlr_pedido", "count"), valor=("vlr_pedido", "sum"))
        .sort_values("qtd", ascending=True)
        .reset_index()
    )

    col_l, col_r = Inches(0.3), Inches(6.8)
    w, h = Inches(6.2), Inches(5.2)

    fig1 = _fig_barras_h(
        ranking["analista"].tolist(),
        ranking["qtd"].tolist(),
        "Quantidade de Liberações por Analista", AZUL_V
    )
    _inserir_grafico(slide, fig1, col_l, Inches(1.5), w, h)

    fig2 = _fig_barras_h(
        ranking["analista"].tolist(),
        ranking["valor"].tolist(),
        "Valor Total Liberado por Analista (R$)", VERDE
    )
    _inserir_grafico(slide, fig2, col_r, Inches(1.5), w, h)


# ============================================================
# SLIDES — REANÁLISES INDIVIDUAL
# ============================================================

def _slide_kpis_reanalises(prs, dados: dict, periodo: str):
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "RESUMO — REANÁLISES DE LIMITE", periodo)

    df = dados.get("df_limites")
    import pandas as pd
    if df is None or (hasattr(df, "empty") and df.empty):
        _texto_sem_dados(slide)
        return
    df = pd.DataFrame(df) if isinstance(df, list) else df

    tot = len(df)
    lim_ant = float(df["limite_anterior"].sum()) if "limite_anterior" in df.columns else 0
    lim_novo = float(df["novo_limite"].sum()) if "novo_limite" in df.columns else 0
    variacao = float(df["variacao"].sum()) if "variacao" in df.columns else 0
    analistas = df["analista"].nunique() if "analista" in df.columns else 0
    maior = float(df["variacao"].max()) if "variacao" in df.columns else 0

    kpis = [
        ("Total Reanálises", str(tot), "no período", AZUL_V),
        ("Limite Anterior", _brl(lim_ant), "soma total carteira", AMARELO),
        ("Novo Limite", _brl(lim_novo), "após reanálises", VERDE),
        ("Variação Total", _brl(variacao), "aumento concedido", VERDE),
        ("Analistas Ativos", str(analistas), "realizaram reanálises", AZUL_V),
        ("Maior Aumento", _brl(maior), "maior concessão individual", AMARELO),
    ]
    _grid_kpis(slide, kpis, Inches(1.5))


def _slide_reanalises_graficos(prs, dados: dict, periodo: str):
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "REANÁLISES — VISÃO CONSOLIDADA", periodo)

    df = dados.get("df_limites")
    import pandas as pd
    if df is None or (hasattr(df, "empty") and df.empty):
        _texto_sem_dados(slide)
        return
    df = pd.DataFrame(df) if isinstance(df, list) else df
    df["variacao"] = pd.to_numeric(df.get("variacao", 0), errors="coerce").fillna(0)

    col_l, col_r = Inches(0.3), Inches(6.8)
    w, h = Inches(6.2), Inches(5.2)

    lim_ant = float(df["limite_anterior"].sum()) if "limite_anterior" in df.columns else 0
    lim_novo = float(df["novo_limite"].sum()) if "novo_limite" in df.columns else 0
    variacao = float(df["variacao"].sum())
    fig1 = _fig_barras_v_cores(
        ["Limite Anterior", "Novo Limite", "Variação"],
        [lim_ant, lim_novo, variacao],
        "Consolidado de Limites (R$)",
        [AMARELO, VERDE, AZUL_V]
    )
    _inserir_grafico(slide, fig1, col_l, Inches(1.5), w, h)

    if "analista" in df.columns:
        por_an = (
            df.groupby("analista")["variacao"].sum()
            .reset_index()
            .sort_values("variacao", ascending=True)
        )
        fig2 = _fig_barras_h(
            por_an["analista"].tolist(),
            por_an["variacao"].tolist(),
            "Variação de Limite Concedida por Analista (R$)", VERDE
        )
        _inserir_grafico(slide, fig2, col_r, Inches(1.5), w, h)


def _slide_top_reanalises(prs, dados: dict, periodo: str):
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "TOP REANÁLISES — MAIORES AUMENTOS", periodo)

    df = dados.get("df_limites")
    import pandas as pd
    if df is None or (hasattr(df, "empty") and df.empty):
        _texto_sem_dados(slide)
        return
    df = pd.DataFrame(df) if isinstance(df, list) else df
    df["variacao"] = pd.to_numeric(df.get("variacao", 0), errors="coerce").fillna(0)
    top = df.nlargest(10, "variacao")

    headers = ["Nome / Parceiro", "Lim. Anterior", "Novo Limite", "Variação", "Analista"]
    cols = ["nome", "limite_anterior", "novo_limite", "variacao", "analista"]
    larguras = [Inches(3.5), Inches(1.8), Inches(1.8), Inches(1.8), Inches(2.0)]
    _tabela_dados(slide, top, headers, cols, larguras, Inches(1.5),
                  formatos={1: "R$", 2: "R$", 3: "R$"})


# ============================================================
# SLIDES — PPT GERAL (ANÁLISE PROFUNDA)
# ============================================================

def _slide_sumario_executivo(prs, dados_cob: list, dados_cred: list):
    from pptx.util import Inches, Pt
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "SUMÁRIO EXECUTIVO", "Consolidado de Todos os Períodos")

    meses_cob = len(dados_cob)
    meses_cred = len(dados_cred)
    total_acordos = sum(d.get("acordos", {}).get("total_acordos", 0) for d in dados_cob)
    total_recebido = sum(d.get("baixas", {}).get("total_recebido", 0) for d in dados_cob)
    total_pedidos = sum(d.get("liberacoes", {}).get("total", 0) for d in dados_cred)
    total_reanalises = sum(d.get("limites", {}).get("total", 0) for d in dados_cred)
    vlr_aprovado = sum(d.get("liberacoes", {}).get("vlr_aprovado", 0) for d in dados_cred)
    variacao_lim = sum(d.get("limites", {}).get("total_variacao", 0) for d in dados_cred)

    kpis = [
        ("Meses Analisados", f"{meses_cob} / {meses_cred}", "Cobrança / Crédito", AZUL_V),
        ("Total Acordos", str(total_acordos), "realizados no período", VERDE),
        ("Total Recebido", _brl(total_recebido), "cobrança acumulada", VERDE),
        ("Total Pedidos Crédito", f"{total_pedidos:,}".replace(",", "."), "processados", AZUL_V),
        ("Valor Aprovado Crédito", _brl(vlr_aprovado), "direto + analista", VERDE),
        ("Limites Concedidos", _brl(variacao_lim), f"{total_reanalises} reanálises", AMARELO),
    ]
    _grid_kpis(slide, kpis, Inches(1.5))


def _slide_evolucao_cobranca_geral(prs, dados_cob: list):
    if not dados_cob:
        return
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "EVOLUÇÃO DE COBRANÇA — MÊS A MÊS", "Histórico Completo")

    meses = [d.get("mes_label", "") for d in dados_cob]
    valores = [d.get("acordos", {}).get("valor_total", 0) for d in dados_cob]
    recebidos = [d.get("baixas", {}).get("total_recebido", 0) for d in dados_cob]
    cancelamentos = [d.get("acordos", {}).get("pct_cancelamento", 0) for d in dados_cob]
    qtd_acordos = [d.get("acordos", {}).get("total_acordos", 0) for d in dados_cob]

    col_l, col_r = Inches(0.3), Inches(6.8)
    w, h = Inches(6.2), Inches(2.4)

    fig1 = _fig_barras_duplas_linha(
        meses, valores, recebidos,
        "Valor Acordado", "Total Recebido",
        "Valor Acordado vs Total Recebido por Mês (R$)"
    )
    _inserir_grafico(slide, fig1, col_l, Inches(1.5), Inches(13.0), h)

    fig2 = _fig_linha_simples(meses, qtd_acordos, "Quantidade de Acordos por Mês", AZUL_V)
    _inserir_grafico(slide, fig2, col_l, Inches(4.1), w, h)

    cores_canc = [VERDE if v <= 10 else VERMELHO for v in cancelamentos]
    fig3 = _fig_barras_v_cores(
        meses, cancelamentos,
        "Taxa de Cancelamento por Mês (%)", cores_canc, linha_ref=10.0
    )
    _inserir_grafico(slide, fig3, col_r, Inches(4.1), w, h)


def _slide_analise_cobradores_geral(prs, dados_cob: list):
    if not dados_cob:
        return
    from pptx.util import Inches
    from collections import defaultdict
    import pandas as pd

    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "DESEMPENHO POR COBRADOR — ACUMULADO", "Histórico Completo")

    por_cob = defaultdict(lambda: {"valor_acordado": 0, "recebido": 0, "acordos": 0,
                                    "tto_soma": 0, "tto_n": 0})
    for d in dados_cob:
        df_ac = d.get("df_acordos")
        df_bx = d.get("df_baixas")
        df_pf = d.get("df_performance")

        if df_ac is not None:
            df_ac = pd.DataFrame(df_ac) if isinstance(df_ac, list) else df_ac
            if "negociador" in df_ac.columns and "valor_total" in df_ac.columns:
                for cob, g in df_ac.groupby("negociador"):
                    por_cob[cob]["valor_acordado"] += g["valor_total"].sum()
                    por_cob[cob]["acordos"] += len(g)

        if df_bx is not None:
            df_bx = pd.DataFrame(df_bx) if isinstance(df_bx, list) else df_bx
            if "cobrador" in df_bx.columns and "vlr_liquido" in df_bx.columns:
                for cob, g in df_bx.groupby("cobrador"):
                    por_cob[cob]["recebido"] += g["vlr_liquido"].sum()

        if df_pf is not None:
            df_pf = pd.DataFrame(df_pf) if isinstance(df_pf, list) else df_pf
            if "cobrador" in df_pf.columns and "pct_aderencia" in df_pf.columns:
                for _, row in df_pf.iterrows():
                    cob = row.get("cobrador", "")
                    if cob:
                        por_cob[cob]["tto_soma"] += float(row.get("pct_aderencia", 0) or 0)
                        por_cob[cob]["tto_n"] += 1

    if not por_cob:
        _texto_sem_dados(slide)
        return

    cobradores = list(por_cob.keys())
    val_ac = [por_cob[c]["valor_acordado"] for c in cobradores]
    receb = [por_cob[c]["recebido"] for c in cobradores]
    ader = [
        round(por_cob[c]["tto_soma"] / por_cob[c]["tto_n"] * 100, 1)
        if por_cob[c]["tto_n"] > 0 else 0
        for c in cobradores
    ]

    col_l, col_r = Inches(0.3), Inches(6.8)
    w, h = Inches(6.2), Inches(2.5)

    fig1 = _fig_barras_duplas(
        cobradores, val_ac, receb,
        "Valor Acordado", "Total Recebido",
        "Acordado vs Recebido por Cobrador (Acumulado R$)"
    )
    _inserir_grafico(slide, fig1, col_l, Inches(1.5), Inches(13.0), h)

    cores_a = [VERDE if v >= 80 else (AMARELO if v >= 60 else VERMELHO) for v in ader]
    fig2 = _fig_barras_v_cores(cobradores, ader, "Aderência Média à Meta (%)", cores_a, linha_ref=100.0)
    _inserir_grafico(slide, fig2, col_l, Inches(4.2), w, h)

    acordos = [por_cob[c]["acordos"] for c in cobradores]
    fig3 = _fig_barras_v_cores(cobradores, acordos, "Total de Acordos Realizados", [AZUL_V]*len(cobradores))
    _inserir_grafico(slide, fig3, col_r, Inches(4.2), w, h)


def _slide_evolucao_credito_geral(prs, dados_cred: list):
    if not dados_cred:
        return
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "EVOLUÇÃO DE CRÉDITO — MÊS A MÊS", "Histórico Completo")

    meses = [d.get("mes_label", "") for d in dados_cred]
    pct_direto = [d.get("liberacoes", {}).get("pct_direto", 0) for d in dados_cred]
    pct_lib = [d.get("liberacoes", {}).get("pct_liberados", 0) for d in dados_cred]
    pct_neg = [d.get("liberacoes", {}).get("pct_negados", 0) for d in dados_cred]
    reanalises = [d.get("limites", {}).get("total", 0) for d in dados_cred]
    variacoes = [d.get("limites", {}).get("total_variacao", 0) for d in dados_cred]

    col_l, col_r = Inches(0.3), Inches(6.8)
    w, h = Inches(6.2), Inches(2.4)

    fig1 = _fig_linhas_multiplas(
        meses,
        [pct_direto, pct_lib, pct_neg],
        ["% Direto", "% Liberado (Analista)", "% Negado"],
        [VERDE, AZUL_V, VERMELHO],
        "Tendência dos Índices de Crédito (%)"
    )
    _inserir_grafico(slide, fig1, col_l, Inches(1.5), Inches(13.0), h)

    fig2 = _fig_barras_v_cores(meses, reanalises, "Quantidade de Reanálises por Mês", [AMARELO]*len(meses))
    _inserir_grafico(slide, fig2, col_l, Inches(4.1), w, h)

    fig3 = _fig_barras_v_cores(meses, variacoes, "Variação de Limite Concedida por Mês (R$)", [VERDE]*len(meses))
    _inserir_grafico(slide, fig3, col_r, Inches(4.1), w, h)


def _slide_analise_analistas_geral(prs, dados_cred: list):
    if not dados_cred:
        return
    from pptx.util import Inches
    from collections import defaultdict
    import pandas as pd

    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "DESEMPENHO POR ANALISTA — ACUMULADO", "Histórico Completo")

    por_an = defaultdict(lambda: {"qtd": 0, "valor": 0, "reanalises": 0, "var_limite": 0})
    for d in dados_cred:
        df_lib = d.get("df_liberacoes")
        df_lim = d.get("df_limites")

        if df_lib is not None:
            df_lib = pd.DataFrame(df_lib) if isinstance(df_lib, list) else df_lib
            df_lib_f = df_lib[df_lib["tipo"] == "LIBERADO"] if "tipo" in df_lib.columns else df_lib
            if "analista" in df_lib_f.columns:
                for an, g in df_lib_f.dropna(subset=["analista"]).groupby("analista"):
                    por_an[an]["qtd"] += len(g)
                    por_an[an]["valor"] += float(g["vlr_pedido"].sum()) if "vlr_pedido" in g.columns else 0

        if df_lim is not None:
            df_lim = pd.DataFrame(df_lim) if isinstance(df_lim, list) else df_lim
            if "analista" in df_lim.columns:
                for an, g in df_lim.dropna(subset=["analista"]).groupby("analista"):
                    por_an[an]["reanalises"] += len(g)
                    por_an[an]["var_limite"] += float(g["variacao"].sum()) if "variacao" in g.columns else 0

    if not por_an:
        _texto_sem_dados(slide)
        return

    analistas = list(por_an.keys())
    qtds = [por_an[a]["qtd"] for a in analistas]
    valores = [por_an[a]["valor"] for a in analistas]
    reanalises = [por_an[a]["reanalises"] for a in analistas]
    var_limite = [por_an[a]["var_limite"] for a in analistas]

    col_l, col_r = Inches(0.3), Inches(6.8)
    w, h = Inches(6.2), Inches(2.5)

    fig1 = _fig_barras_duplas(
        analistas, qtds, reanalises,
        "Liberações", "Reanálises",
        "Liberações e Reanálises por Analista (Acumulado)"
    )
    _inserir_grafico(slide, fig1, col_l, Inches(1.5), Inches(13.0), h)

    fig2 = _fig_barras_h(analistas, valores, "Valor Total Liberado por Analista (R$)", AZUL_V)
    _inserir_grafico(slide, fig2, col_l, Inches(4.2), w, h)

    fig3 = _fig_barras_h(analistas, var_limite, "Variação de Limite Concedida por Analista (R$)", VERDE)
    _inserir_grafico(slide, fig3, col_r, Inches(4.2), w, h)


def _slide_tendencias_geral(prs, dados_cob: list, dados_cred: list):
    from pptx.util import Inches, Pt
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "TENDÊNCIAS E CORRELAÇÕES", "Visão Integrada")

    meses_cob = [d.get("mes_label", "") for d in dados_cob]
    meses_cred = [d.get("mes_label", "") for d in dados_cred]

    col_l, col_r = Inches(0.3), Inches(6.8)
    w, h = Inches(6.2), Inches(2.4)

    if dados_cob:
        aderencias = []
        for d in dados_cob:
            pf = d.get("performance", {})
            a = pf.get("aderencia_media_pct", 0) or 0
            aderencias.append(a)
        cores_ad = [VERDE if v >= 80 else (AMARELO if v >= 60 else VERMELHO) for v in aderencias]
        fig1 = _fig_barras_v_cores(meses_cob, aderencias, "Aderência Média dos Cobradores por Mês (%)",
                                    cores_ad, linha_ref=80.0)
        _inserir_grafico(slide, fig1, col_l, Inches(1.5), w, h)

    if dados_cred:
        pct_negados = [d.get("liberacoes", {}).get("pct_negados", 0) for d in dados_cred]
        cores_neg = [VERDE if v <= 1.5 else (AMARELO if v <= 3 else VERMELHO) for v in pct_negados]
        fig2 = _fig_barras_v_cores(meses_cred, pct_negados, "Taxa de Negação por Mês (%)",
                                    cores_neg, linha_ref=2.0)
        _inserir_grafico(slide, fig2, col_r, Inches(1.5), w, h)

    if dados_cob:
        cancelamentos = [d.get("acordos", {}).get("pct_cancelamento", 0) for d in dados_cob]
        recebidos = [d.get("baixas", {}).get("total_recebido", 0) for d in dados_cob]
        fig3 = _fig_barras_duplas_linha(
            meses_cob, recebidos, cancelamentos,
            "Recebido (R$)", "% Cancelamento",
            "Recebimento vs Cancelamento por Mês"
        )
        _inserir_grafico(slide, fig3, col_l, Inches(4.1), Inches(13.0), h)


def _slide_analise_profunda_texto(prs, dados_cob: list, dados_cred: list):
    from pptx.util import Inches, Pt
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "ANÁLISE GLOBAL — CRÉDITO & COBRANÇA", "Observações e Conclusões")

    linhas = []

    # --- Cobrança ---
    if dados_cob:
        total_acordos = sum(d.get("acordos", {}).get("total_acordos", 0) for d in dados_cob)
        total_recebido = sum(d.get("baixas", {}).get("total_recebido", 0) for d in dados_cob)
        media_ader = sum(d.get("performance", {}).get("aderencia_media_pct", 0) or 0 for d in dados_cob) / len(dados_cob)
        media_canc = sum(d.get("acordos", {}).get("pct_cancelamento", 0) for d in dados_cob) / len(dados_cob)

        linhas.append(("COBRANÇA — SÍNTESE", AZUL_V, True))
        linhas.append((f"• {total_acordos} acordos realizados no período, totalizando {_brl(total_recebido)} recebidos.", AZUL, False))
        linhas.append((f"• Aderência média dos cobradores à meta de 2h/dia: {_pct(media_ader)}.", AZUL, False))

        if media_ader < 80:
            linhas.append((f"  ⚠ Aderência abaixo de 80% — recomenda-se acompanhamento mais próximo do TTO.", VERMELHO, False))
        else:
            linhas.append((f"  ✓ Equipe manteve boa disciplina de atendimento ao longo do período.", VERDE, False))

        if media_canc > 10:
            linhas.append((f"• Taxa média de cancelamento de {_pct(media_canc)} — acima do ideal. Avaliar causas.", VERMELHO, False))
        else:
            linhas.append((f"• Taxa de cancelamento controlada: {_pct(media_canc)} em média.", VERDE, False))

        linhas.append(("", BRANCO, False))

    # --- Crédito ---
    if dados_cred:
        total_ped = sum(d.get("liberacoes", {}).get("total", 0) for d in dados_cred)
        media_dir = sum(d.get("liberacoes", {}).get("pct_direto", 0) for d in dados_cred) / len(dados_cred)
        media_neg = sum(d.get("liberacoes", {}).get("pct_negados", 0) for d in dados_cred) / len(dados_cred)
        total_reanalises = sum(d.get("limites", {}).get("total", 0) for d in dados_cred)
        total_var = sum(d.get("limites", {}).get("total_variacao", 0) for d in dados_cred)

        linhas.append(("CRÉDITO — SÍNTESE", AZUL_V, True))
        linhas.append((f"• {total_ped:,} pedidos processados no período.".replace(",", "."), AZUL, False))
        linhas.append((f"• {_pct(media_dir)} dos pedidos passaram direto (automático) em média — alto grau de automação.", VERDE if media_dir >= 70 else AMARELO, False))
        linhas.append((f"• Taxa média de negação: {_pct(media_neg)}." + (" Dentro do aceitável." if media_neg <= 2 else " Requer atenção."), VERDE if media_neg <= 2 else VERMELHO, False))
        linhas.append((f"• {total_reanalises} reanálises de limite realizadas, com {_brl(total_var)} de crédito adicional concedido.", AZUL, False))

    _bloco_texto_analise(slide, linhas)


def _slide_alertas_recomendacoes(prs, dados_cob: list, dados_cred: list):
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, "ALERTAS E RECOMENDAÇÕES", "Pontos de Atenção")

    alertas = []
    positivos = []

    for d in dados_cob:
        mes = d.get("mes_label", "")
        pct_c = d.get("acordos", {}).get("pct_cancelamento", 0)
        ader = d.get("performance", {}).get("aderencia_media_pct", 0) or 0

        if pct_c > 15:
            alertas.append(f"⚠ {mes}: Cancelamentos em {_pct(pct_c)} — investigar motivos")
        if ader < 60:
            alertas.append(f"⚠ {mes}: Aderência crítica em {_pct(ader)} — TTO muito baixo")
        if pct_c <= 10 and ader >= 80:
            positivos.append(f"✓ {mes} (Cobrança): Bom desempenho geral")

    for d in dados_cred:
        mes = d.get("mes_label", "")
        pct_neg = d.get("liberacoes", {}).get("pct_negados", 0)
        pct_dir = d.get("liberacoes", {}).get("pct_direto", 0)

        if pct_neg > 3:
            alertas.append(f"⚠ {mes}: Negações em {_pct(pct_neg)} — revisar critérios de aprovação")
        if pct_dir < 65:
            alertas.append(f"⚠ {mes}: Automação baixa {_pct(pct_dir)} — carga manual elevada")
        if pct_neg <= 1.5 and pct_dir >= 75:
            positivos.append(f"✓ {mes} (Crédito): Alta automação e baixa rejeição")

    if not alertas and not positivos:
        alertas.append("Nenhum alerta crítico identificado no período.")

    linhas = []
    if alertas:
        linhas.append(("ALERTAS", VERMELHO, True))
        for a in alertas[:6]:
            linhas.append((a, "5C0000" if "⚠" in a else AZUL, False))
        linhas.append(("", BRANCO, False))
    if positivos:
        linhas.append(("DESTAQUES POSITIVOS", VERDE, True))
        for p in positivos[:5]:
            linhas.append((p, "1A4D1F", False))

    _bloco_texto_analise(slide, linhas)


# ============================================================
# SLIDE DE TEXTO — ANÁLISE TEXTUAL
# ============================================================

def _slide_destaques_texto(prs, linhas: list, titulo: str, periodo: str):
    from pptx.util import Inches
    slide = _novo_slide(prs)
    _cabecalho(slide, prs, titulo, periodo)
    _bloco_texto_analise(slide, linhas)


def _analise_cobranca(dados: dict) -> list:
    ac = dados.get("acordos", {})
    bx = dados.get("baixas", {})
    pf = dados.get("performance", {})
    linhas = []
    linhas.append(("ACORDOS", AZUL_V, True))
    linhas.append((f"• {ac.get('total_acordos', 0)} acordos realizados · Valor total: {_brl(ac.get('valor_total', 0))}", AZUL, False))
    linhas.append((f"• Ticket médio por acordo: {_brl(ac.get('ticket_medio', 0))}", AZUL, False))
    pct_c = ac.get('pct_cancelamento', 0)
    linhas.append((f"• Taxa de cancelamento: {_pct(pct_c)}" + (" — atenção!" if pct_c > 10 else " — controlada"), VERMELHO if pct_c > 10 else VERDE, False))
    linhas.append(("", BRANCO, False))
    linhas.append(("COBRANÇA BAIXADA", AZUL_V, True))
    linhas.append((f"• Total recebido: {_brl(bx.get('total_recebido', 0))} em {bx.get('total_baixas', 0)} baixas", AZUL, False))
    linhas.append(("", BRANCO, False))
    linhas.append(("PERFORMANCE", AZUL_V, True))
    ader = pf.get("aderencia_media_pct", 0) or 0
    linhas.append((f"• Aderência média à meta: {_pct(ader)}" + (" ✓" if ader >= 80 else " — abaixo do ideal"), VERDE if ader >= 80 else VERMELHO, False))
    top = pf.get("top_cobrador", "")
    if top:
        linhas.append((f"• Melhor cobrador do período: {top}", VERDE, False))
    return linhas


def _analise_credito(dados: dict) -> list:
    lib = dados.get("liberacoes", {})
    lim = dados.get("limites", {})
    linhas = []
    linhas.append(("LIBERAÇÕES POR PEDIDO", AZUL_V, True))
    linhas.append((f"• {lib.get('total', 0):,} pedidos processados no período".replace(",", "."), AZUL, False))
    linhas.append((f"• Direto (automático): {_pct(lib.get('pct_direto', 0))} ({lib.get('qtd_direto', 0):,} pedidos)".replace(",", "."), VERDE, False))
    linhas.append((f"• Liberado (analista): {_pct(lib.get('pct_liberados', 0))} ({lib.get('qtd_liberados', 0):,} pedidos)".replace(",", "."), AZUL, False))
    pct_neg = lib.get("pct_negados", 0)
    linhas.append((f"• Negados: {_pct(pct_neg)} ({lib.get('qtd_negados', 0):,} pedidos)".replace(",", ".") + (" — atenção!" if pct_neg > 2 else ""), VERMELHO if pct_neg > 2 else AZUL, False))
    linhas.append((f"• Valor aprovado total: {_brl(lib.get('vlr_aprovado', 0))}", VERDE, False))
    linhas.append(("", BRANCO, False))
    linhas.append(("REANÁLISES DE LIMITE", AZUL_V, True))
    linhas.append((f"• {lim.get('total', 0)} reanálises · {_brl(lim.get('total_variacao', 0))} de crédito adicional concedido", AZUL, False))
    return linhas


def _analise_reanalises(dados: dict) -> list:
    import pandas as pd
    df = dados.get("df_limites")
    if df is None or (hasattr(df, "empty") and df.empty):
        return [("Sem dados suficientes para análise.", CINZA_M, False)]
    df = pd.DataFrame(df) if isinstance(df, list) else df
    df["variacao"] = pd.to_numeric(df.get("variacao", 0), errors="coerce").fillna(0)

    tot = len(df)
    lim_ant = float(df["limite_anterior"].sum()) if "limite_anterior" in df.columns else 0
    lim_novo = float(df["novo_limite"].sum()) if "novo_limite" in df.columns else 0
    variacao = float(df["variacao"].sum())
    pct_var = (variacao / lim_ant * 100) if lim_ant > 0 else 0
    analistas = df["analista"].nunique() if "analista" in df.columns else 0
    maior = float(df["variacao"].max())
    nome_maior = df.loc[df["variacao"].idxmax(), "nome"] if "nome" in df.columns else "—"

    linhas = []
    linhas.append(("RESUMO DAS REANÁLISES", AZUL_V, True))
    linhas.append((f"• {tot} reanálises realizadas por {analistas} analista(s)", AZUL, False))
    linhas.append((f"• Limite anterior total: {_brl(lim_ant)} → Novo limite: {_brl(lim_novo)}", AZUL, False))
    linhas.append((f"• Variação total concedida: {_brl(variacao)} (+{_pct(pct_var)} sobre o limite anterior)", VERDE, False))
    linhas.append((f"• Maior concessão individual: {_brl(maior)} — {nome_maior}", VERDE, False))
    linhas.append(("", BRANCO, False))
    linhas.append(("OBSERVAÇÕES", AZUL_V, True))
    if pct_var > 50:
        linhas.append(("• Crescimento expressivo de limites — monitorar exposição de crédito", VERMELHO, False))
    elif pct_var > 20:
        linhas.append(("• Expansão moderada de crédito — dentro de parâmetros normais", AMARELO, False))
    else:
        linhas.append(("• Concessão conservadora — baixo impacto na exposição total", VERDE, False))
    return linhas


# ============================================================
# PRIMITIVAS DE SLIDE
# ============================================================

def _novo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.util import Inches
    # Fundo branco
    _ret(slide, 0, 0, prs.slide_width, prs.slide_height, BRANCO)
    # Faixa azul topo
    _ret(slide, 0, 0, prs.slide_width, Inches(1.35), AZUL)
    # Linha amarela
    _ret(slide, 0, Inches(1.33), prs.slide_width, Inches(0.045), AMARELO)
    # Borda lateral esquerda
    _ret(slide, 0, 0, Inches(0.07), prs.slide_height, AMARELO)
    return slide


def _cabecalho(slide, prs, titulo: str, subtitulo: str):
    from pptx.util import Inches, Pt
    _txt(slide, titulo, Inches(0.2), Inches(0.1), Inches(10.5), Inches(0.7),
         Pt(18), BRANCO, bold=True)
    _txt(slide, subtitulo, Inches(0.2), Inches(0.8), Inches(9), Inches(0.38),
         Pt(10), AMARELO)
    _txt(slide, "GRUPO LLE", Inches(10.8), Inches(0.22), Inches(2.3), Inches(0.4),
         Pt(11), AMARELO, bold=True)


def _slide_capa(prs, titulo: str, periodo: str, subtitulo: str):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = prs.slide_width, prs.slide_height
    _ret(slide, 0, 0, w, h, AZUL)
    _ret(slide, 0, 0, Inches(0.5), h, AMARELO)
    _ret(slide, 0, h - Inches(0.07), w, Inches(0.07), AMARELO)
    _txt(slide, titulo.upper(), Inches(1), Inches(2.0), Inches(11), Inches(1.4),
         Pt(40), BRANCO, bold=True)
    _ret(slide, Inches(1), Inches(3.55), Inches(5), Inches(0.045), AMARELO)
    _txt(slide, periodo, Inches(1), Inches(3.7), Inches(9), Inches(0.6),
         Pt(22), AMARELO, bold=True)
    _txt(slide, subtitulo, Inches(1), Inches(4.45), Inches(10), Inches(0.45),
         Pt(13), "AAAAAA")
    _txt(slide, "GRUPO LLE  ·  Gestão Financeira", Inches(1), Inches(6.75),
         Inches(10), Inches(0.4), Pt(10), AMARELO)


def _slide_encerramento(prs, area: str):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = prs.slide_width, prs.slide_height
    _ret(slide, 0, 0, w, h, AZUL)
    _ret(slide, 0, 0, Inches(0.5), h, AMARELO)
    _ret(slide, 0, h - Inches(0.07), w, Inches(0.07), AMARELO)
    _txt(slide, "Obrigado.", Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
         Pt(52), BRANCO, bold=True)
    _txt(slide, f"Equipe de {area.title()}  ·  Grupo LLE Ferragens",
         Inches(1.5), Inches(4.1), Inches(10), Inches(0.5), Pt(16), AMARELO)
    _txt(slide, "GRUPO LLE  ·  Gestão Financeira",
         Inches(1.5), Inches(6.6), Inches(10), Inches(0.4), Pt(11), "AAAAAA")


def _texto_sem_dados(slide):
    from pptx.util import Inches, Pt
    _txt(slide, "Sem dados disponíveis para este período.",
         Inches(1), Inches(3.5), Inches(11), Inches(0.5), Pt(14), CINZA_M)


def _grid_kpis(slide, kpis: list, top):
    from pptx.util import Inches, Pt, Emu
    cols = 3
    card_w = Inches(3.95)
    card_h = Inches(1.52)
    gap_x = Inches(0.22)
    gap_y = Inches(0.18)
    start_x = Inches(0.25)

    for i, (label, valor, sub, cor) in enumerate(kpis):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = top + row * (card_h + gap_y)

        card = slide.shapes.add_shape(1, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(CINZA_F)
        card.line.color.rgb = _rgb(cor)
        card.line.width = Emu(30000)

        barra = slide.shapes.add_shape(1, x, y, card_w, Inches(0.065))
        barra.fill.solid()
        barra.fill.fore_color.rgb = _rgb(cor)
        barra.line.fill.background()

        _txt(slide, label, x + Inches(0.14), y + Inches(0.1),
             card_w - Inches(0.28), Inches(0.32), Pt(10), CINZA_M)
        _txt(slide, valor, x + Inches(0.14), y + Inches(0.45),
             card_w - Inches(0.28), Inches(0.65), Pt(24), cor, bold=True)
        _txt(slide, sub, x + Inches(0.14), y + Inches(1.18),
             card_w - Inches(0.28), Inches(0.28), Pt(9), CINZA_M)


def _bloco_texto_analise(slide, linhas: list):
    from pptx.util import Inches, Pt
    y = Inches(1.52)
    for texto, cor, negrito in linhas:
        if not texto:
            y += Inches(0.1)
            continue
        h = Inches(0.3) if negrito else Inches(0.28)
        if negrito:
            # Fundo da seção
            fundo = slide.shapes.add_shape(1, Inches(0.25), y - Inches(0.02),
                                           Inches(12.9), Inches(0.32))
            fundo.fill.solid()
            fundo.fill.fore_color.rgb = _rgb(CINZA_L)
            fundo.line.fill.background()
        _txt(slide, texto, Inches(0.35), y, Inches(12.6), h,
             Pt(11) if negrito else Pt(10), cor, bold=negrito)
        y += h + Inches(0.06)


def _tabela_dados(slide, df, headers, cols, larguras, top, formatos=None):
    from pptx.util import Inches, Pt
    import pandas as pd
    formatos = formatos or {}
    x_start = Inches(0.25)
    y = top
    h_row = Inches(0.32)

    # Cabeçalho
    x = x_start
    for i, (h, w) in enumerate(zip(headers, larguras)):
        cell = slide.shapes.add_shape(1, x, y, w, h_row)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(AZUL)
        cell.line.fill.background()
        _txt(slide, h, x + Inches(0.05), y + Inches(0.04),
             w - Inches(0.1), Inches(0.24), Pt(9), BRANCO, bold=True)
        x += w
    y += h_row + Inches(0.02)

    for idx, (_, row) in enumerate(df.iterrows()):
        if y > Inches(7.0):
            break
        bg = CINZA_F if idx % 2 == 0 else BRANCO
        x = x_start
        for ci, (col, w) in enumerate(zip(cols, larguras)):
            val = row.get(col)
            if pd.isna(val) if not isinstance(val, str) else False:
                val = None
            # Formata
            fmt = formatos.get(ci)
            if fmt == "R$" and val is not None:
                try:
                    val = _brl(float(val))
                except Exception:
                    pass
            elif fmt == "%" and val is not None:
                try:
                    v = float(val)
                    val = f"{v*100:.1f}%" if v <= 1 else f"{v:.1f}%"
                except Exception:
                    pass
            elif fmt == "h" and val is not None:
                try:
                    val = f"{float(val):.2f}h"
                except Exception:
                    pass

            cell = slide.shapes.add_shape(1, x, y, w, h_row)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(bg)
            cell.line.color.rgb = _rgb("DDDDDD")
            _txt(slide, str(val) if val is not None else "—",
                 x + Inches(0.05), y + Inches(0.04),
                 w - Inches(0.1), Inches(0.24), Pt(9), AZUL)
            x += w
        y += h_row


# ============================================================
# GRÁFICOS — MATPLOTLIB
# ============================================================

def _inserir_grafico(slide, fig, left, top, width, height):
    img_bytes = _fig_para_png(fig)
    if img_bytes:
        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width, height)


def _fig_para_png(fig) -> Optional[bytes]:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                    facecolor="white", edgecolor="none")
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception:
        return None


def _base_fig(titulo: str, figsize=(9, 3.8)):
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        fig, ax = plt.subplots(figsize=figsize)
        fig.patch.set_facecolor("white")
        ax.set_facecolor("white")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color("#DDD")
        ax.spines["bottom"].set_color("#DDD")
        ax.tick_params(labelsize=8, colors="#555")
        if titulo:
            ax.set_title(titulo, fontsize=10, fontweight="bold", pad=8, color=f"#{AZUL}")
        return fig, ax
    except Exception:
        return None, None


def _fig_barras_h(labels, valores, titulo, cor):
    fig, ax = _base_fig(titulo)
    if ax is None:
        return None
    import matplotlib.pyplot as plt
    cores = [f"#{cor}"] * len(labels)
    bars = ax.barh(labels, valores, color=cores, edgecolor="white", linewidth=0.5)
    for bar, val in zip(bars, valores):
        if val and val > 0:
            ax.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height()/2,
                    _brl(val), va='center', ha='left', fontsize=7.5, color=f"#{AZUL}")
    ax.xaxis.grid(True, color="#EEE", linewidth=0.5)
    ax.set_axisbelow(True)
    plt.tight_layout(pad=1.2)
    return fig


def _fig_barras_v_cores(labels, valores, titulo, cores, linha_ref=None):
    fig, ax = _base_fig(titulo)
    if ax is None:
        return None
    import matplotlib.pyplot as plt
    cores_hex = [f"#{c}" if not str(c).startswith("#") else c for c in cores]
    bars = ax.bar(labels, valores, color=cores_hex, edgecolor="white", linewidth=0.5)
    for bar, val in zip(bars, valores):
        if val and val > 0:
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() * 1.01,
                    _fmt_val_simples(val), ha='center', va='bottom', fontsize=7.5, color=f"#{AZUL}")
    if linha_ref is not None:
        ax.axhline(y=linha_ref, color=f"#{AMARELO}", linestyle="--", linewidth=1.2, label=f"Ref. {linha_ref}")
    ax.yaxis.grid(True, color="#EEE", linewidth=0.5)
    ax.set_axisbelow(True)
    plt.xticks(rotation=15 if len(labels) > 4 else 0, ha='right' if len(labels) > 4 else 'center')
    plt.tight_layout(pad=1.2)
    return fig


def _fig_barras_v_cores_linha(labels, valores, cores, linha_ref, titulo):
    fig, ax = _base_fig(titulo)
    if ax is None:
        return None
    import matplotlib.pyplot as plt
    cores_hex = [f"#{c}" if not str(c).startswith("#") else c for c in cores]
    bars = ax.bar(labels, valores, color=cores_hex, edgecolor="white", linewidth=0.5)
    for bar, val in zip(bars, valores):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.02,
                f"{val:.2f}h", ha='center', va='bottom', fontsize=8, color=f"#{AZUL}")
    ax.axhline(y=linha_ref, color=f"#{AMARELO}", linestyle="--", linewidth=1.5,
               label=f"Meta {linha_ref}h")
    ax.legend(fontsize=8)
    ax.yaxis.grid(True, color="#EEE", linewidth=0.5)
    ax.set_axisbelow(True)
    ax.set_ylabel("Horas", fontsize=8)
    plt.tight_layout(pad=1.2)
    return fig


def _fig_pizza(labels, valores, titulo):
    fig, ax = _base_fig(titulo, figsize=(6, 4))
    if ax is None:
        return None
    import matplotlib.pyplot as plt
    cores = [f"#{VERDE}", f"#{AZUL_V}", f"#{VERMELHO}", f"#{AMARELO}", f"#{AZUL}"]
    total = sum(valores)
    valores_f = [v for v in valores if v > 0]
    labels_f = [labels[i] for i, v in enumerate(valores) if v > 0]
    if not valores_f:
        ax.text(0.5, 0.5, "Sem dados", ha="center", va="center", transform=ax.transAxes)
        plt.tight_layout()
        return fig
    wedges, texts, autotexts = ax.pie(
        valores_f, labels=labels_f,
        colors=cores[:len(valores_f)],
        autopct=lambda p: f'{p:.1f}%\n({int(round(p*total/100)):,})'.replace(",", "."),
        startangle=90,
        wedgeprops=dict(edgecolor='white', linewidth=1.5),
        textprops={"fontsize": 8.5},
    )
    for autotext in autotexts:
        autotext.set_fontsize(7.5)
    ax.axis('equal')
    plt.tight_layout(pad=1.2)
    return fig


def _fig_barras_duplas(labels, val1, val2, label1, label2, titulo):
    fig, ax = _base_fig(titulo)
    if ax is None:
        return None
    import matplotlib.pyplot as plt
    import numpy as np
    x = np.arange(len(labels))
    w = 0.38
    ax.bar(x - w/2, val1, w, label=label1, color=f"#{AZUL_V}", edgecolor="white")
    ax.bar(x + w/2, val2, w, label=label2, color=f"#{VERDE}", edgecolor="white")
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=8, rotation=15 if len(labels) > 4 else 0,
                       ha='right' if len(labels) > 4 else 'center')
    ax.legend(fontsize=8)
    ax.yaxis.grid(True, color="#EEE", linewidth=0.5)
    ax.set_axisbelow(True)
    plt.tight_layout(pad=1.2)
    return fig


def _fig_barras_duplas_linha(labels, val1, val2, label1, label2, titulo):
    fig, ax = _base_fig(titulo)
    if ax is None:
        return None
    import matplotlib.pyplot as plt
    ax2 = ax.twinx()
    ax.bar(labels, val1, color=f"#{AZUL_V}", alpha=0.8, label=label1, edgecolor="white")
    ax2.plot(labels, val2, color=f"#{VERDE}", marker="o", linewidth=2, markersize=5, label=label2)
    ax.set_ylabel(label1, fontsize=8, color=f"#{AZUL_V}")
    ax2.set_ylabel(label2, fontsize=8, color=f"#{VERDE}")
    ax.tick_params(axis='x', rotation=15 if len(labels) > 4 else 0, labelsize=8)
    lines1, labels1 = ax.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax.legend(lines1 + lines2, labels1 + labels2, fontsize=8, loc="upper left")
    ax.yaxis.grid(True, color="#EEE", linewidth=0.5)
    plt.tight_layout(pad=1.2)
    return fig


def _fig_linha_simples(labels, valores, titulo, cor):
    fig, ax = _base_fig(titulo)
    if ax is None:
        return None
    import matplotlib.pyplot as plt
    ax.plot(labels, valores, color=f"#{cor}", marker="o", linewidth=2, markersize=6)
    for i, (x, y) in enumerate(zip(labels, valores)):
        ax.text(i, y + max(valores)*0.02, str(y), ha='center', va='bottom', fontsize=8)
    ax.yaxis.grid(True, color="#EEE", linewidth=0.5)
    ax.set_axisbelow(True)
    ax.tick_params(axis='x', rotation=15 if len(labels) > 4 else 0, labelsize=8)
    ax.fill_between(range(len(labels)), valores, alpha=0.1, color=f"#{cor}")
    plt.tight_layout(pad=1.2)
    return fig


def _fig_linhas_multiplas(labels, series, nomes, cores, titulo):
    fig, ax = _base_fig(titulo)
    if ax is None:
        return None
    import matplotlib.pyplot as plt
    for vals, nome, cor in zip(series, nomes, cores):
        ax.plot(labels, vals, color=f"#{cor}", marker="o", linewidth=2,
                markersize=5, label=nome)
    ax.legend(fontsize=8)
    ax.yaxis.grid(True, color="#EEE", linewidth=0.5)
    ax.set_axisbelow(True)
    ax.set_ylabel("%", fontsize=8)
    ax.tick_params(axis='x', rotation=15 if len(labels) > 4 else 0, labelsize=8)
    plt.tight_layout(pad=1.2)
    return fig


# ============================================================
# HELPERS
# ============================================================

def _nova_apresentacao():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def _salvar(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _rgb(hex6: str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex6)


def _ret(slide, x, y, w, h, cor):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(cor)
    shape.line.fill.background()
    return shape


def _txt(slide, texto: str, x, y, w, h, tamanho, cor: str, bold: bool = False):
    from pptx.enum.text import PP_ALIGN
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(texto)
    p.font.size = tamanho
    p.font.bold = bold
    p.font.color.rgb = _rgb(cor)
    try:
        p.font.name = "Calibri"
    except Exception:
        pass
    return tx


def _brl(v) -> str:
    try:
        f = float(v)
        return f"R$ {f:,.0f}".replace(",", "X").replace(".", ",").replace("X", ".")
    except Exception:
        return "R$ 0"


def _pct(v) -> str:
    try:
        f = float(v)
        return f"{f:.1f}%"
    except Exception:
        return "0%"


def _fmt_val_simples(v) -> str:
    try:
        f = float(v)
        if f >= 1_000_000:
            return f"R${f/1_000_000:.1f}M"
        if f >= 1_000:
            return f"R${f/1_000:.0f}K"
        return f"{f:.1f}"
    except Exception:
        return str(v)
